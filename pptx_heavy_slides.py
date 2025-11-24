#!/Users/julian/Programming/Python/JNF_Experiments/slide_size/venv/bin/python
"""
PowerPoint Heavy Slides Analyzer

Analyzes .pptx files to identify which slides contribute most to file size
due to embedded media (images, videos, audio).
"""

import argparse
import csv
import io
import json
import logging
import sys
from pathlib import Path
from typing import TypedDict
from collections import defaultdict

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


__version__ = "1.0.0"


# Data model
class MediaItem(TypedDict):
    """Represents a single media item on a slide."""
    type: str  # "image", "video", "audio", "other"
    size_bytes: int
    filename: str | None
    content_type: str | None
    relationship_id: str | None
    shared: bool


class SlideMediaStats(TypedDict):
    """Statistics about media content on a single slide."""
    slide_index: int  # 1-based index
    slide_title: str | None
    total_media_bytes: int
    image_bytes: int
    video_bytes: int
    audio_bytes: int
    other_media_bytes: int
    media_items: list[MediaItem]


class ImageDimensions(TypedDict):
    """Dimensions of an image in pixels and display size."""
    pixel_width: int
    pixel_height: int
    display_width_px: int
    display_height_px: int
    resolution_ratio: float  # How many times larger the image is vs display


class OptimizationOpportunity(TypedDict):
    """Represents a potential optimization for an image."""
    slide_index: int
    slide_title: str | None
    opportunity_type: str  # "oversized_resolution", "absolute_size", "png_photo", "uncompressed_jpeg"
    current_bytes: int
    potential_bytes: int
    savings_bytes: int
    savings_percent: float
    current_dimensions: str  # e.g., "3840x2160"
    display_dimensions: str  # e.g., "384x216"
    recommended_dimensions: str  # e.g., "768x432"
    current_format: str
    recommended_format: str
    details: str
    severity: str  # "high", "medium", "low"
    is_shared: bool


class LayoutMediaStats(TypedDict):
    """Statistics about media content in a slide layout."""
    layout_name: str
    layout_index: int  # Index within master
    total_media_bytes: int
    image_bytes: int
    video_bytes: int
    audio_bytes: int
    other_media_bytes: int
    media_count: int
    is_used: bool  # Whether any slide uses this layout
    slides_using: list[int]  # 1-based slide indices using this layout


class MasterMediaStats(TypedDict):
    """Statistics about media content in a slide master."""
    master_index: int  # 1-based
    master_name: str | None
    total_media_bytes: int  # Media directly on master
    image_bytes: int
    video_bytes: int
    audio_bytes: int
    other_media_bytes: int
    media_count: int
    layouts: list[LayoutMediaStats]
    total_layout_bytes: int  # Sum of all layout media
    unused_layout_bytes: int  # Media in unused layouts


class MastersReport(TypedDict):
    """Complete report on slide masters and layouts."""
    total_masters: int
    total_layouts: int
    unused_layouts: int
    total_master_media_bytes: int
    total_layout_media_bytes: int
    unused_layout_media_bytes: int
    masters: list[MasterMediaStats]


def setup_logging(verbose: bool = False) -> None:
    """Configure logging based on verbosity level."""
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format='%(levelname)s: %(message)s'
    )


def get_slide_title(slide) -> str | None:
    """
    Extract the title from a slide.

    Args:
        slide: A slide object from python-pptx

    Returns:
        The slide title text, or None if no title found
    """
    if not slide.shapes.title:
        return None

    try:
        title_text = slide.shapes.title.text.strip()
        return title_text if title_text else None
    except (AttributeError, IndexError):
        return None


def get_image_dimensions(image_blob: bytes, shape) -> ImageDimensions:
    """
    Extract pixel and display dimensions for an image.

    Args:
        image_blob: The raw image bytes
        shape: The shape object from python-pptx

    Returns:
        ImageDimensions with pixel and display information
    """
    # Get pixel dimensions using Pillow
    img = Image.open(io.BytesIO(image_blob))
    pixel_width, pixel_height = img.size

    # Convert shape dimensions from EMUs to pixels (assuming 96 DPI)
    # 1 inch = 914400 EMUs, 96 DPI means 96 pixels per inch
    display_width_px = int(shape.width / 914400 * 96)
    display_height_px = int(shape.height / 914400 * 96)

    # Calculate resolution ratio (how much larger the image is than display)
    if display_width_px > 0 and display_height_px > 0:
        width_ratio = pixel_width / display_width_px
        height_ratio = pixel_height / display_height_px
        resolution_ratio = max(width_ratio, height_ratio)
    else:
        resolution_ratio = 1.0

    return ImageDimensions(
        pixel_width=pixel_width,
        pixel_height=pixel_height,
        display_width_px=display_width_px,
        display_height_px=display_height_px,
        resolution_ratio=resolution_ratio
    )


def analyze_image_optimization(
    image_blob: bytes,
    content_type: str,
    dimensions: ImageDimensions,
    slide_index: int,
    slide_title: str | None,
    is_shared: bool
) -> list[OptimizationOpportunity]:
    """
    Analyze an image for optimization opportunities.

    Conservative thresholds for conference-quality presentations:
    - Allows 2x resolution for retina displays
    - Flags truly wasteful oversizing (>2.5x display size)
    - Identifies format and quality improvements

    Args:
        image_blob: Raw image bytes
        content_type: MIME type (e.g., 'image/jpeg')
        dimensions: Image dimensions info
        slide_index: Slide number (1-based)
        slide_title: Slide title or None
        is_shared: Whether image is shared across slides

    Returns:
        List of optimization opportunities found
    """
    opportunities = []
    current_bytes = len(image_blob)
    img = Image.open(io.BytesIO(image_blob))
    img_format = img.format

    current_dim_str = f"{dimensions['pixel_width']}x{dimensions['pixel_height']}"
    display_dim_str = f"{dimensions['display_width_px']}x{dimensions['display_height_px']}"

    # 1. Check for oversized resolution (>2.5x display size)
    # This is the most common and impactful issue
    if dimensions['resolution_ratio'] > 2.5:
        # Recommend 2x for retina quality
        recommended_width = dimensions['display_width_px'] * 2
        recommended_height = dimensions['display_height_px'] * 2

        # Maintain aspect ratio
        aspect_ratio = dimensions['pixel_width'] / dimensions['pixel_height']
        if recommended_width / recommended_height > aspect_ratio:
            recommended_width = int(recommended_height * aspect_ratio)
        else:
            recommended_height = int(recommended_width / aspect_ratio)

        recommended_dim_str = f"{recommended_width}x{recommended_height}"

        # Estimate savings: proportional to pixel reduction
        pixel_reduction = (recommended_width * recommended_height) / (
            dimensions['pixel_width'] * dimensions['pixel_height']
        )
        potential_bytes = int(current_bytes * pixel_reduction)
        savings_bytes = current_bytes - potential_bytes

        severity = "high" if dimensions['resolution_ratio'] > 5 else "medium"

        opportunities.append(OptimizationOpportunity(
            slide_index=slide_index,
            slide_title=slide_title,
            opportunity_type="oversized_resolution",
            current_bytes=current_bytes,
            potential_bytes=potential_bytes,
            savings_bytes=savings_bytes,
            savings_percent=round((savings_bytes / current_bytes) * 100, 1),
            current_dimensions=current_dim_str,
            display_dimensions=display_dim_str,
            recommended_dimensions=recommended_dim_str,
            current_format=img_format or content_type,
            recommended_format=img_format or content_type,
            details=f"Image is {dimensions['resolution_ratio']:.1f}x larger than display size. "
                   f"Resizing to 2x (retina quality) would maintain sharpness on all screens.",
            severity=severity,
            is_shared=is_shared
        ))

    # 2. Check for absolute size caps (>3200px on longest edge)
    # Safety net for unreasonably large images
    max_dimension = max(dimensions['pixel_width'], dimensions['pixel_height'])
    if max_dimension > 3200:
        # Recommend 2560px max (covers retina 1280px displays, suitable for conference projectors)
        target_max = 2560
        aspect_ratio = dimensions['pixel_width'] / dimensions['pixel_height']

        if dimensions['pixel_width'] > dimensions['pixel_height']:
            recommended_width = target_max
            recommended_height = int(target_max / aspect_ratio)
        else:
            recommended_height = target_max
            recommended_width = int(target_max * aspect_ratio)

        recommended_dim_str = f"{recommended_width}x{recommended_height}"

        pixel_reduction = (recommended_width * recommended_height) / (
            dimensions['pixel_width'] * dimensions['pixel_height']
        )
        potential_bytes = int(current_bytes * pixel_reduction)
        savings_bytes = current_bytes - potential_bytes

        # Only add if not already caught by oversized resolution check
        if not any(opp['opportunity_type'] == 'oversized_resolution' for opp in opportunities):
            opportunities.append(OptimizationOpportunity(
                slide_index=slide_index,
                slide_title=slide_title,
                opportunity_type="absolute_size",
                current_bytes=current_bytes,
                potential_bytes=potential_bytes,
                savings_bytes=savings_bytes,
                savings_percent=round((savings_bytes / current_bytes) * 100, 1),
                current_dimensions=current_dim_str,
                display_dimensions=display_dim_str,
                recommended_dimensions=recommended_dim_str,
                current_format=img_format or content_type,
                recommended_format=img_format or content_type,
                details=f"Image exceeds {max_dimension}px. Conference projectors rarely exceed "
                       f"1920x1080 (Full HD). Recommend max {target_max}px for high-quality projection.",
                severity="medium",
                is_shared=is_shared
            ))

    # 3. Check for PNG photos (should be JPEG)
    # PNG is great for screenshots/diagrams, wasteful for photos
    if img_format == 'PNG' and current_bytes > 1_000_000:  # >1MB
        # Estimate JPEG size at quality 85 (roughly 70-80% reduction for photos)
        potential_bytes = int(current_bytes * 0.3)  # Conservative estimate
        savings_bytes = current_bytes - potential_bytes

        opportunities.append(OptimizationOpportunity(
            slide_index=slide_index,
            slide_title=slide_title,
            opportunity_type="png_photo",
            current_bytes=current_bytes,
            potential_bytes=potential_bytes,
            savings_bytes=savings_bytes,
            savings_percent=round((savings_bytes / current_bytes) * 100, 1),
            current_dimensions=current_dim_str,
            display_dimensions=display_dim_str,
            recommended_dimensions=current_dim_str,
            current_format="PNG",
            recommended_format="JPEG",
            details=f"Large PNG file ({format_bytes(current_bytes)}). Converting to JPEG "
                   f"at quality 85-90 provides visually identical results for photos, "
                   f"with significant file size reduction.",
            severity="medium" if current_bytes > 3_000_000 else "low",
            is_shared=is_shared
        ))

    # 4. Check for uncompressed/high-quality JPEG
    # JPEG with >1 byte/pixel suggests quality 95-100 (often unnecessary for presentations)
    if img_format == 'JPEG':
        bytes_per_pixel = current_bytes / (dimensions['pixel_width'] * dimensions['pixel_height'])
        if bytes_per_pixel > 1.0:
            # Estimate re-saving at quality 85 (roughly 50% reduction)
            potential_bytes = int(current_bytes * 0.5)
            savings_bytes = current_bytes - potential_bytes

            opportunities.append(OptimizationOpportunity(
                slide_index=slide_index,
                slide_title=slide_title,
                opportunity_type="uncompressed_jpeg",
                current_bytes=current_bytes,
                potential_bytes=potential_bytes,
                savings_bytes=savings_bytes,
                savings_percent=round((savings_bytes / current_bytes) * 100, 1),
                current_dimensions=current_dim_str,
                display_dimensions=display_dim_str,
                recommended_dimensions=current_dim_str,
                current_format="JPEG (high quality)",
                recommended_format="JPEG (quality 85)",
                details=f"JPEG file appears to use very high compression quality "
                       f"({bytes_per_pixel:.2f} bytes/pixel). Re-saving at quality 85 "
                       f"produces visually identical results for conference projection.",
                severity="low",
                is_shared=is_shared
            ))

    return opportunities


def analyze_pptx_media(path: str, include_shared_media: bool = False) -> list[SlideMediaStats]:
    """
    Analyze a PowerPoint file to determine media size per slide.

    Args:
        path: Path to the .pptx file
        include_shared_media: If True, count shared media on every slide.
                             If False (default), count shared media only on first appearance.

    Returns:
        List of SlideMediaStats, one per slide, sorted by total_media_bytes descending

    Raises:
        FileNotFoundError: If the file doesn't exist
        ValueError: If the file is not a .pptx file or is corrupt
    """
    # Validate file
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"file not found: {path}")

    if file_path.suffix.lower() != '.pptx':
        raise ValueError(f"unsupported file type (expected .pptx): {path}")

    # Try to open the presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"failed to open .pptx: {e}")

    logging.info(f"Found {len(prs.slides)} slides")

    # Track all media across slides to detect sharing
    media_registry = {}  # key: image blob hash or media part name -> value: {size, slides, first_slide}
    slide_media_map = defaultdict(list)  # slide_index -> list of media items

    # First pass: collect all media and track where it appears
    for slide_idx, slide in enumerate(prs.slides, start=1):
        logging.debug(f"Analyzing slide {slide_idx}")

        for shape in slide.shapes:
            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    size_bytes = len(blob)
                    content_type = image.content_type

                    # Use blob hash as identifier for detecting duplicates
                    media_key = hash(blob)

                    if media_key not in media_registry:
                        media_registry[media_key] = {
                            'size': size_bytes,
                            'type': 'image',
                            'content_type': content_type,
                            'filename': image.filename if hasattr(image, 'filename') else None,
                            'slides': [],
                            'first_slide': slide_idx
                        }

                    media_registry[media_key]['slides'].append(slide_idx)

                    # Store reference for this slide
                    slide_media_map[slide_idx].append({
                        'media_key': media_key,
                        'type': 'image',
                        'size_bytes': size_bytes,
                        'content_type': content_type,
                        'filename': media_registry[media_key]['filename']
                    })

                    logging.debug(f"  Found image: {size_bytes} bytes, type={content_type}")

                except Exception as e:
                    logging.warning(f"  Failed to extract image from slide {slide_idx}: {e}")

            # Handle video and audio
            elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                try:
                    media_format = shape.media_format
                    if media_format:
                        # Get the media part
                        media_part = shape.part.related_part(shape._element.xpath('.//a:videoFile | .//a:audioFile | .//p:media')[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'))

                        size_bytes = len(media_part.blob)
                        content_type = media_part.content_type

                        # Determine media type
                        if 'video' in content_type.lower():
                            media_type = 'video'
                        elif 'audio' in content_type.lower():
                            media_type = 'audio'
                        else:
                            media_type = 'other'

                        # Use part name as identifier
                        media_key = media_part.partname

                        if media_key not in media_registry:
                            media_registry[media_key] = {
                                'size': size_bytes,
                                'type': media_type,
                                'content_type': content_type,
                                'filename': Path(media_part.partname).name,
                                'slides': [],
                                'first_slide': slide_idx
                            }

                        media_registry[media_key]['slides'].append(slide_idx)

                        slide_media_map[slide_idx].append({
                            'media_key': media_key,
                            'type': media_type,
                            'size_bytes': size_bytes,
                            'content_type': content_type,
                            'filename': media_registry[media_key]['filename']
                        })

                        logging.debug(f"  Found {media_type}: {size_bytes} bytes, type={content_type}")

                except Exception as e:
                    logging.warning(f"  Failed to extract media from slide {slide_idx}: {e}")

    # Mark shared media
    for media_key, info in media_registry.items():
        if len(info['slides']) > 1:
            logging.debug(f"Media {info.get('filename', media_key)} appears on slides: {info['slides']}")

    # Second pass: build SlideMediaStats for each slide
    results = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = get_slide_title(slide)
        media_items = []

        image_bytes = 0
        video_bytes = 0
        audio_bytes = 0
        other_media_bytes = 0

        # Process media for this slide
        for media_ref in slide_media_map.get(slide_idx, []):
            media_key = media_ref['media_key']
            media_info = media_registry[media_key]

            is_shared = len(media_info['slides']) > 1
            is_first_appearance = media_info['first_slide'] == slide_idx

            # Determine if we should count the bytes
            if include_shared_media or not is_shared or is_first_appearance:
                count_bytes = media_ref['size_bytes']
            else:
                count_bytes = 0  # Shared media, not first appearance, don't count

            # Create media item
            media_item: MediaItem = {
                'type': media_ref['type'],
                'size_bytes': count_bytes,
                'filename': media_ref['filename'],
                'content_type': media_ref['content_type'],
                'relationship_id': None,  # Could be enhanced later
                'shared': is_shared
            }
            media_items.append(media_item)

            # Accumulate by type
            if media_ref['type'] == 'image':
                image_bytes += count_bytes
            elif media_ref['type'] == 'video':
                video_bytes += count_bytes
            elif media_ref['type'] == 'audio':
                audio_bytes += count_bytes
            else:
                other_media_bytes += count_bytes

        total_media_bytes = image_bytes + video_bytes + audio_bytes + other_media_bytes

        stats: SlideMediaStats = {
            'slide_index': slide_idx,
            'slide_title': title,
            'total_media_bytes': total_media_bytes,
            'image_bytes': image_bytes,
            'video_bytes': video_bytes,
            'audio_bytes': audio_bytes,
            'other_media_bytes': other_media_bytes,
            'media_items': media_items
        }
        results.append(stats)

    # Sort by total_media_bytes descending
    results.sort(key=lambda x: x['total_media_bytes'], reverse=True)

    return results


def format_bytes(num_bytes: int) -> str:
    """Format bytes into human-readable string (KB, MB, GB)."""
    if num_bytes == 0:
        return "0.0 MB"

    for unit in ['B', 'KB', 'MB', 'GB']:
        if num_bytes < 1024.0 or unit == 'GB':
            if unit == 'B':
                return f"{num_bytes} {unit}"
            return f"{num_bytes:.1f} {unit}"
        num_bytes /= 1024.0

    return f"{num_bytes:.1f} GB"


def print_console_output(results: list[SlideMediaStats], filename: str, top_n: int | None = None) -> None:
    """
    Print human-readable console output.

    Args:
        results: List of SlideMediaStats (should be pre-sorted)
        filename: Name of the analyzed file
        top_n: If specified, show only top N slides
    """
    print(f"\nAnalyzing: {filename}")
    print(f"\nTotal slides: {len(results)}")
    print("\nRanked by media size (descending):\n")

    display_results = results[:top_n] if top_n else results

    for rank, stats in enumerate(display_results, start=1):
        slide_num = stats['slide_index']
        size_str = format_bytes(stats['total_media_bytes'])
        title = stats['slide_title'] if stats['slide_title'] else "(no title)"

        print(f"#{rank:<3} Slide {slide_num:<3} | {size_str:>10} | title=\"{title}\"")

    print()


def write_json_output(results: list[SlideMediaStats], output_path: str) -> None:
    """
    Write results to JSON file.

    Args:
        results: List of SlideMediaStats
        output_path: Path to output file

    Raises:
        IOError: If file write fails
    """
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        logging.info(f"JSON output written to: {output_path}")
    except Exception as e:
        raise IOError(f"failed to write output file {output_path}: {e}")


def write_csv_output(results: list[SlideMediaStats], output_path: str) -> None:
    """
    Write results to CSV file.

    Args:
        results: List of SlideMediaStats
        output_path: Path to output file

    Raises:
        IOError: If file write fails
    """
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            # Header
            writer.writerow([
                'rank',
                'slide_index',
                'slide_title',
                'total_media_bytes',
                'image_bytes',
                'video_bytes',
                'audio_bytes',
                'other_media_bytes'
            ])

            # Data rows
            for rank, stats in enumerate(results, start=1):
                writer.writerow([
                    rank,
                    stats['slide_index'],
                    stats['slide_title'] or '',
                    stats['total_media_bytes'],
                    stats['image_bytes'],
                    stats['video_bytes'],
                    stats['audio_bytes'],
                    stats['other_media_bytes']
                ])

        logging.info(f"CSV output written to: {output_path}")
    except Exception as e:
        raise IOError(f"failed to write output file {output_path}: {e}")


def analyze_optimization_opportunities(path: str) -> list[OptimizationOpportunity]:
    """
    Analyze a PowerPoint file for image optimization opportunities.

    Args:
        path: Path to the .pptx file

    Returns:
        List of OptimizationOpportunity, sorted by potential savings descending

    Raises:
        FileNotFoundError: If the file doesn't exist
        ValueError: If the file is not a .pptx file or is corrupt
    """
    # Validate file
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"file not found: {path}")

    if file_path.suffix.lower() != '.pptx':
        raise ValueError(f"unsupported file type (expected .pptx): {path}")

    # Open presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"failed to open .pptx: {e}")

    logging.info(f"Analyzing {len(prs.slides)} slides for optimization opportunities")

    # Track shared media
    media_registry = {}  # key: image blob hash -> value: {size, slides, first_slide}
    slide_shapes_map = {}  # slide_index -> list of (shape, media_key)

    # First pass: collect all images and detect sharing
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_shapes_map[slide_idx] = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    media_key = hash(blob)

                    if media_key not in media_registry:
                        media_registry[media_key] = {
                            'blob': blob,
                            'content_type': image.content_type,
                            'slides': [],
                            'first_slide': slide_idx
                        }

                    media_registry[media_key]['slides'].append(slide_idx)
                    slide_shapes_map[slide_idx].append((shape, media_key))

                except Exception as e:
                    logging.warning(f"Failed to analyze image on slide {slide_idx}: {e}")

    # Second pass: analyze each unique image for optimization opportunities
    all_opportunities = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = get_slide_title(slide)

        for shape, media_key in slide_shapes_map.get(slide_idx, []):
            media_info = media_registry[media_key]
            is_shared = len(media_info['slides']) > 1
            is_first_appearance = media_info['first_slide'] == slide_idx

            # Only analyze each unique image once (on first appearance)
            if not is_first_appearance:
                continue

            try:
                # Get image dimensions
                dimensions = get_image_dimensions(media_info['blob'], shape)

                # Analyze for optimization opportunities
                opportunities = analyze_image_optimization(
                    image_blob=media_info['blob'],
                    content_type=media_info['content_type'],
                    dimensions=dimensions,
                    slide_index=slide_idx,
                    slide_title=slide_title,
                    is_shared=is_shared
                )

                all_opportunities.extend(opportunities)

            except Exception as e:
                logging.warning(f"Failed to analyze optimization for slide {slide_idx}: {e}")

    # Sort by potential savings (highest first)
    all_opportunities.sort(key=lambda x: x['savings_bytes'], reverse=True)

    return all_opportunities


def print_optimization_report(opportunities: list[OptimizationOpportunity], filename: str) -> None:
    """
    Print human-readable optimization report to console.

    Args:
        opportunities: List of optimization opportunities
        filename: Name of the analyzed file
    """
    if not opportunities:
        print(f"\nOptimization Report: {filename}")
        print("\nNo optimization opportunities found. Your images are well-optimized!")
        return

    # Calculate total potential savings
    total_savings = sum(opp['savings_bytes'] for opp in opportunities)
    total_current = sum(opp['current_bytes'] for opp in opportunities)
    total_savings_pct = (total_savings / total_current * 100) if total_current > 0 else 0

    # Group by severity
    high_severity = [opp for opp in opportunities if opp['severity'] == 'high']
    medium_severity = [opp for opp in opportunities if opp['severity'] == 'medium']
    low_severity = [opp for opp in opportunities if opp['severity'] == 'low']

    print(f"\n{'='*80}")
    print(f"OPTIMIZATION REPORT: {filename}")
    print(f"{'='*80}")
    print(f"\nSUMMARY:")
    print(f"  Total opportunities found: {len(opportunities)}")
    print(f"  Potential savings: {format_bytes(total_savings)} ({total_savings_pct:.1f}% reduction)")
    print(f"  High priority: {len(high_severity)} | Medium: {len(medium_severity)} | Low: {len(low_severity)}")

    print(f"\n{'='*80}")
    print(f"RECOMMENDATIONS (sorted by potential savings):")
    print(f"{'='*80}\n")

    for idx, opp in enumerate(opportunities, start=1):
        severity_marker = {
            'high': '🔴 HIGH',
            'medium': '🟡 MEDIUM',
            'low': '🟢 LOW'
        }.get(opp['severity'], opp['severity'])

        print(f"#{idx} - Slide {opp['slide_index']}: {opp['slide_title'] or '(no title)'}")
        print(f"    Priority: {severity_marker}")
        print(f"    Current: {format_bytes(opp['current_bytes'])} | {opp['current_format']} | {opp['current_dimensions']}")
        print(f"    Display size: {opp['display_dimensions']} pixels")
        print(f"    Recommended: {opp['recommended_format']} | {opp['recommended_dimensions']}")
        print(f"    Potential savings: {format_bytes(opp['savings_bytes'])} ({opp['savings_percent']}%)")

        if opp['is_shared']:
            print(f"    ⚠️  SHARED: This image appears on multiple slides - optimization affects all")

        print(f"    💡 {opp['details']}")
        print()

    print(f"{'='*80}")
    print(f"NOTES FOR CONFERENCE PRESENTATIONS:")
    print(f"{'='*80}")
    print(f"  • Most conference projectors are 1920x1080 (Full HD)")
    print(f"  • 2x resolution (e.g., 1536x864 for 768x432 display) ensures retina quality")
    print(f"  • Images larger than 2560px rarely improve visual quality on projectors")
    print(f"  • JPEG quality 85-90 is visually identical to quality 95-100 when projected")
    print(f"  • PNG is best for screenshots/diagrams; JPEG is best for photos")
    print()


def analyze_slide_masters(path: str) -> MastersReport:
    """
    Analyze slide masters and layouts for media content and usage.

    Args:
        path: Path to the .pptx file

    Returns:
        MastersReport with complete analysis of masters and layouts

    Raises:
        FileNotFoundError: If the file doesn't exist
        ValueError: If the file is not a .pptx file or is corrupt
    """
    # Validate file
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"file not found: {path}")

    if file_path.suffix.lower() != '.pptx':
        raise ValueError(f"unsupported file type (expected .pptx): {path}")

    # Open presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"failed to open .pptx: {e}")

    logging.info(f"Analyzing {len(prs.slide_masters)} slide masters")

    # Build a map of which layouts are used by which slides
    layout_usage = {}  # layout id -> list of slide indices
    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_id = id(slide.slide_layout)
        if layout_id not in layout_usage:
            layout_usage[layout_id] = []
        layout_usage[layout_id].append(slide_idx)

    masters_stats = []
    total_master_media = 0
    total_layout_media = 0
    total_unused_layout_media = 0
    total_layouts = 0
    unused_layouts = 0

    # Analyze each slide master
    for master_idx, master in enumerate(prs.slide_masters, start=1):
        # Get master name
        master_name = None
        try:
            # Try to get theme name or other identifier
            if hasattr(master, 'name') and master.name:
                master_name = master.name
        except Exception:
            pass

        # Analyze media directly on the master
        master_image_bytes = 0
        master_video_bytes = 0
        master_audio_bytes = 0
        master_other_bytes = 0
        master_media_count = 0

        for shape in master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    master_image_bytes += len(blob)
                    master_media_count += 1
                except Exception as e:
                    logging.warning(f"Failed to extract image from master {master_idx}: {e}")

        master_total = master_image_bytes + master_video_bytes + master_audio_bytes + master_other_bytes
        total_master_media += master_total

        # Analyze each layout in this master
        layout_stats = []
        master_layout_bytes = 0
        master_unused_bytes = 0

        for layout_idx, layout in enumerate(master.slide_layouts, start=1):
            total_layouts += 1
            layout_id = id(layout)
            slides_using = layout_usage.get(layout_id, [])
            is_used = len(slides_using) > 0

            if not is_used:
                unused_layouts += 1

            # Get layout name
            layout_name = layout.name if hasattr(layout, 'name') and layout.name else f"Layout {layout_idx}"

            # Analyze media in this layout
            layout_image_bytes = 0
            layout_video_bytes = 0
            layout_audio_bytes = 0
            layout_other_bytes = 0
            layout_media_count = 0

            for shape in layout.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        layout_image_bytes += len(blob)
                        layout_media_count += 1
                    except Exception as e:
                        logging.warning(f"Failed to extract image from layout {layout_name}: {e}")

            layout_total = layout_image_bytes + layout_video_bytes + layout_audio_bytes + layout_other_bytes
            master_layout_bytes += layout_total
            total_layout_media += layout_total

            if not is_used:
                master_unused_bytes += layout_total
                total_unused_layout_media += layout_total

            layout_stat: LayoutMediaStats = {
                'layout_name': layout_name,
                'layout_index': layout_idx,
                'total_media_bytes': layout_total,
                'image_bytes': layout_image_bytes,
                'video_bytes': layout_video_bytes,
                'audio_bytes': layout_audio_bytes,
                'other_media_bytes': layout_other_bytes,
                'media_count': layout_media_count,
                'is_used': is_used,
                'slides_using': slides_using
            }
            layout_stats.append(layout_stat)

        master_stat: MasterMediaStats = {
            'master_index': master_idx,
            'master_name': master_name,
            'total_media_bytes': master_total,
            'image_bytes': master_image_bytes,
            'video_bytes': master_video_bytes,
            'audio_bytes': master_audio_bytes,
            'other_media_bytes': master_other_bytes,
            'media_count': master_media_count,
            'layouts': layout_stats,
            'total_layout_bytes': master_layout_bytes,
            'unused_layout_bytes': master_unused_bytes
        }
        masters_stats.append(master_stat)

    report: MastersReport = {
        'total_masters': len(prs.slide_masters),
        'total_layouts': total_layouts,
        'unused_layouts': unused_layouts,
        'total_master_media_bytes': total_master_media,
        'total_layout_media_bytes': total_layout_media,
        'unused_layout_media_bytes': total_unused_layout_media,
        'masters': masters_stats
    }

    return report


def print_masters_report(report: MastersReport, filename: str) -> None:
    """
    Print human-readable masters report to console.

    Args:
        report: MastersReport from analyze_slide_masters
        filename: Name of the analyzed file
    """
    print(f"\n{'='*80}")
    print(f"SLIDE MASTERS REPORT: {filename}")
    print(f"{'='*80}")

    # Summary
    print(f"\nSUMMARY:")
    print(f"  Total slide masters: {report['total_masters']}")
    print(f"  Total layouts: {report['total_layouts']}")
    print(f"  Unused layouts: {report['unused_layouts']}")
    print(f"\n  Media in masters: {format_bytes(report['total_master_media_bytes'])}")
    print(f"  Media in layouts: {format_bytes(report['total_layout_media_bytes'])}")
    print(f"  Media in UNUSED layouts: {format_bytes(report['unused_layout_media_bytes'])}")

    if report['unused_layout_media_bytes'] > 0:
        print(f"\n  ⚠️  You could save {format_bytes(report['unused_layout_media_bytes'])} by deleting unused layouts")

    # Details per master
    for master in report['masters']:
        print(f"\n{'-'*80}")
        master_name = master['master_name'] or f"Master {master['master_index']}"
        print(f"MASTER {master['master_index']}: {master_name}")
        print(f"{'-'*80}")

        if master['media_count'] > 0:
            print(f"  Media on master: {format_bytes(master['total_media_bytes'])} ({master['media_count']} items)")
        else:
            print(f"  Media on master: (none)")

        print(f"  Layouts: {len(master['layouts'])} total, {format_bytes(master['total_layout_bytes'])} media")

        if master['unused_layout_bytes'] > 0:
            print(f"  ⚠️  Unused layout media: {format_bytes(master['unused_layout_bytes'])}")

        # List layouts with media or unused status
        layouts_with_media = [l for l in master['layouts'] if l['total_media_bytes'] > 0]
        unused_layouts = [l for l in master['layouts'] if not l['is_used']]

        if layouts_with_media:
            print(f"\n  Layouts with media:")
            for layout in sorted(layouts_with_media, key=lambda x: x['total_media_bytes'], reverse=True):
                status = "UNUSED" if not layout['is_used'] else f"used by {len(layout['slides_using'])} slides"
                print(f"    • {layout['layout_name']}: {format_bytes(layout['total_media_bytes'])} [{status}]")

        if unused_layouts:
            unused_with_no_media = [l for l in unused_layouts if l['total_media_bytes'] == 0]
            if unused_with_no_media:
                print(f"\n  Unused layouts (no media):")
                for layout in unused_with_no_media:
                    print(f"    • {layout['layout_name']}")

    # Recommendations
    if report['unused_layouts'] > 0:
        print(f"\n{'='*80}")
        print(f"RECOMMENDATIONS:")
        print(f"{'='*80}")
        print(f"\n  To reduce file size, consider deleting unused layouts:")
        print(f"  1. Open the presentation in PowerPoint")
        print(f"  2. Go to View > Slide Master")
        print(f"  3. Right-click unused layouts and select 'Delete Layout'")
        print(f"  4. Close the Slide Master view")

        if report['unused_layout_media_bytes'] > 0:
            print(f"\n  Potential savings: {format_bytes(report['unused_layout_media_bytes'])}")

    print()


def delete_unused_layouts(path: str) -> tuple[int, int, str]:
    """
    Delete unused layouts from a presentation.

    Args:
        path: Path to the .pptx file

    Returns:
        Tuple of (layouts_deleted, bytes_saved, output_path)

    Raises:
        FileNotFoundError: If the file doesn't exist
        ValueError: If the file is not a .pptx file or is corrupt
    """
    # Validate file
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"file not found: {path}")

    if file_path.suffix.lower() != '.pptx':
        raise ValueError(f"unsupported file type (expected .pptx): {path}")

    # Open presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"failed to open .pptx: {e}")

    # Build a set of used layout ids
    used_layout_ids = set()
    for slide in prs.slides:
        used_layout_ids.add(id(slide.slide_layout))

    # Find and delete unused layouts
    layouts_deleted = 0
    bytes_saved = 0

    for master in prs.slide_masters:
        # Collect unused layouts (must iterate in reverse to avoid index issues)
        layouts_to_delete = []
        for layout in master.slide_layouts:
            if id(layout) not in used_layout_ids:
                # Calculate bytes in this layout
                layout_bytes = 0
                for shape in layout.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            layout_bytes += len(shape.image.blob)
                        except Exception:
                            pass
                layouts_to_delete.append((layout, layout_bytes))

        # Delete unused layouts by removing from the sldLayoutIdLst
        sldLayoutIdLst = master._element.get_or_add_sldLayoutIdLst()

        for layout, layout_bytes in layouts_to_delete:
            try:
                # Find and remove the sldLayoutId entry for this layout
                for sldLayoutId in list(sldLayoutIdLst):
                    rId = sldLayoutId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                    if rId:
                        rel = master.part.rels.get(rId)
                        if rel and rel.target_part == layout.part:
                            sldLayoutIdLst.remove(sldLayoutId)
                            layouts_deleted += 1
                            bytes_saved += layout_bytes
                            logging.debug(f"Deleted layout: {layout.name}")
                            break
            except Exception as e:
                logging.warning(f"Failed to delete layout {layout.name}: {e}")

    # Save to new file
    output_path = str(file_path.with_stem(file_path.stem + "_cleaned"))
    try:
        prs.save(output_path)
        logging.info(f"Saved cleaned presentation to: {output_path}")
    except Exception as e:
        raise IOError(f"failed to save cleaned presentation: {e}")

    # Note: Orphaned media files may remain in the package.
    # Use PowerPoint's "Compress Pictures" or File > Info > Compress Media
    # to fully clean up after deleting layouts.

    return layouts_deleted, bytes_saved, output_path


def main() -> int:
    """Main CLI entry point. Returns exit code."""
    # Show description if no arguments provided
    if len(sys.argv) == 1:
        print("PowerPoint Heavy Slides Analyzer")
        print("=" * 35)
        print("\nAnalyzes .pptx files to identify which slides contribute most")
        print("to file size due to embedded media (images, videos, audio).")
        print("\nUsage:")
        print("  pptx-analyzer <file.pptx>                    Analyze media by slide")
        print("  pptx-analyzer <file.pptx> --top 5            Show top 5 heaviest slides")
        print("  pptx-analyzer <file.pptx> --optimization-report")
        print("                                               Find image optimization opportunities")
        print("  pptx-analyzer <file.pptx> --masters-report   Analyze slide masters/layouts")
        print("  pptx-analyzer <file.pptx> --delete-unused-layouts")
        print("                                               Remove unused layouts")
        print("\nRun with --help for all options.")
        return 0

    parser = argparse.ArgumentParser(
        description='Analyze PowerPoint files to identify heavy slides with embedded media.',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )

    parser.add_argument(
        'input_path',
        help='Path to the .pptx file to analyze'
    )

    parser.add_argument(
        '--top',
        type=int,
        metavar='N',
        help='Show only the top N heaviest slides'
    )

    parser.add_argument(
        '--output-json',
        metavar='PATH',
        help='Write results as JSON to the specified path'
    )

    parser.add_argument(
        '--output-csv',
        metavar='PATH',
        help='Write results as CSV to the specified path'
    )

    media_group = parser.add_mutually_exclusive_group()
    media_group.add_argument(
        '--include-shared-media',
        action='store_true',
        help='Count shared media on every slide that uses it'
    )
    media_group.add_argument(
        '--ignore-shared-media',
        action='store_true',
        default=True,
        help='Count shared media only on first slide (default)'
    )

    parser.add_argument(
        '--optimization-report',
        action='store_true',
        help='Generate optimization recommendations for reducing file size (conference-quality focused)'
    )

    parser.add_argument(
        '--masters-report',
        action='store_true',
        help='Generate report on slide masters and layouts, including unused layouts with media'
    )

    parser.add_argument(
        '--delete-unused-layouts',
        action='store_true',
        help='Delete unused layouts and save as new file (original is preserved)'
    )

    parser.add_argument(
        '--verbose',
        action='store_true',
        help='Enable verbose debug logging'
    )

    parser.add_argument(
        '--version',
        action='version',
        version=f'%(prog)s {__version__}'
    )

    args = parser.parse_args()

    # Setup logging
    setup_logging(args.verbose)

    try:
        # Check if delete unused layouts is requested
        if args.delete_unused_layouts:
            # First show what will be deleted
            report = analyze_slide_masters(args.input_path)
            print_masters_report(report, args.input_path)

            if report['unused_layouts'] > 0:
                # Delete unused layouts
                layouts_deleted, bytes_saved, output_path = delete_unused_layouts(args.input_path)
                print(f"\n{'='*80}")
                print(f"CLEANUP COMPLETE")
                print(f"{'='*80}")
                print(f"  Layouts deleted: {layouts_deleted}")
                print(f"  Saved to: {output_path}")
                print(f"  Original file preserved: {args.input_path}")
                if bytes_saved > 0:
                    print(f"\n  To fully reclaim {format_bytes(bytes_saved)} from orphaned media:")
                    print(f"  1. Open {output_path} in PowerPoint")
                    print(f"  2. File > Info > Compress Media (or Compress Pictures)")
                    print(f"  3. Save the file")
                print()
            else:
                print("\nNo unused layouts to delete.")
        # Check if masters report is requested
        elif args.masters_report:
            # Run masters analysis
            report = analyze_slide_masters(args.input_path)
            print_masters_report(report, args.input_path)
        # Check if optimization report is requested
        elif args.optimization_report:
            # Run optimization analysis
            opportunities = analyze_optimization_opportunities(args.input_path)
            print_optimization_report(opportunities, args.input_path)
        else:
            # Run standard media analysis
            results = analyze_pptx_media(
                args.input_path,
                include_shared_media=args.include_shared_media
            )

            # Console output (always show)
            print_console_output(results, args.input_path, args.top)

            # Optional JSON output
            if args.output_json:
                write_json_output(results, args.output_json)

            # Optional CSV output
            if args.output_csv:
                write_csv_output(results, args.output_csv)

        return 0

    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except IOError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except Exception as e:
        print(f"Error: unexpected error: {e}", file=sys.stderr)
        logging.exception("Unexpected error occurred")
        return 1


if __name__ == "__main__":
    sys.exit(main())
