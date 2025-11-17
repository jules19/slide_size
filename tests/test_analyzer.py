"""
Tests for the PowerPoint Heavy Slides Analyzer.
"""

import io
import json
import os
import pytest
import tempfile
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches

# Import the module under test
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
import pptx_heavy_slides


def create_test_image(width: int = 100, height: int = 100, color: str = 'red') -> bytes:
    """Create a simple test image and return as bytes."""
    img = Image.new('RGB', (width, height), color=color)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    return img_bytes.getvalue()


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


def test_single_slide_with_image(temp_dir):
    """Test analysis of a single-slide deck with one image."""
    # Create a presentation with one slide and one image
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Test Slide"

    # Add an image
    img_bytes = create_test_image(200, 200, 'blue')
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(3))

    # Save presentation
    pptx_path = temp_dir / "single_image.pptx"
    prs.save(str(pptx_path))

    # Analyze
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))

    # Assertions
    assert len(results) == 1
    assert results[0]['slide_index'] == 1
    assert results[0]['slide_title'] == "Test Slide"
    assert results[0]['total_media_bytes'] > 0
    assert results[0]['image_bytes'] > 0
    assert results[0]['video_bytes'] == 0
    assert results[0]['audio_bytes'] == 0
    assert len(results[0]['media_items']) == 1
    assert results[0]['media_items'][0]['type'] == 'image'
    assert results[0]['media_items'][0]['shared'] is False


def test_shared_media_ignore_mode(temp_dir):
    """Test that shared media is counted only on first slide by default."""
    # Create a presentation with same image on multiple slides
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create a single image to share
    img_bytes = create_test_image(300, 300, 'green')

    # Add 3 slides with the same image
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = f"Slide {i+1}"

        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(2))

    # Save presentation
    pptx_path = temp_dir / "shared_image.pptx"
    prs.save(str(pptx_path))

    # Analyze with default (ignore shared media)
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path), include_shared_media=False)

    assert len(results) == 3

    # Find which slide got counted (should be slide 1, the first appearance)
    slide_with_bytes = [r for r in results if r['total_media_bytes'] > 0]
    slides_without_bytes = [r for r in results if r['total_media_bytes'] == 0]

    assert len(slide_with_bytes) == 1, "Only one slide should have media bytes counted"
    assert len(slides_without_bytes) == 2, "Two slides should have 0 bytes"

    # The first slide should have the bytes
    assert slide_with_bytes[0]['slide_index'] == 1

    # All media items should be marked as shared
    for result in results:
        for item in result['media_items']:
            assert item['shared'] is True


def test_shared_media_include_mode(temp_dir):
    """Test that shared media is counted on every slide with --include-shared-media."""
    # Create a presentation with same image on multiple slides
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    img_bytes = create_test_image(300, 300, 'yellow')

    # Add 3 slides with the same image
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = f"Slide {i+1}"

        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(2))

    # Save presentation
    pptx_path = temp_dir / "shared_image_include.pptx"
    prs.save(str(pptx_path))

    # Analyze with include_shared_media=True
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path), include_shared_media=True)

    assert len(results) == 3

    # All slides should have the same media bytes
    sizes = [r['total_media_bytes'] for r in results]
    assert all(size > 0 for size in sizes), "All slides should have bytes counted"
    assert len(set(sizes)) == 1, "All slides should have the same size"


def test_empty_deck(temp_dir):
    """Test analysis of a deck with no media."""
    # Create a presentation with slides but no media
    prs = Presentation()

    # Add 3 slides with just text
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        title = slide.shapes.title
        if title:
            title.text = f"Slide {i+1}"

    # Save presentation
    pptx_path = temp_dir / "empty.pptx"
    prs.save(str(pptx_path))

    # Analyze
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))

    assert len(results) == 3
    for result in results:
        assert result['total_media_bytes'] == 0
        assert result['image_bytes'] == 0
        assert result['video_bytes'] == 0
        assert result['audio_bytes'] == 0
        assert len(result['media_items']) == 0


def test_multiple_images_different_sizes(temp_dir):
    """Test deck with multiple images of different sizes for proper ranking."""
    prs = Presentation()

    # Slide 1: small image
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    if slide1.shapes.title:
        slide1.shapes.title.text = "Small Image"
    img_small = io.BytesIO(create_test_image(50, 50, 'red'))
    slide1.shapes.add_picture(img_small, Inches(1), Inches(1), width=Inches(1))

    # Slide 2: large image
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    if slide2.shapes.title:
        slide2.shapes.title.text = "Large Image"
    img_large = io.BytesIO(create_test_image(1000, 1000, 'blue'))
    slide2.shapes.add_picture(img_large, Inches(1), Inches(1), width=Inches(4))

    # Slide 3: medium image
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    if slide3.shapes.title:
        slide3.shapes.title.text = "Medium Image"
    img_medium = io.BytesIO(create_test_image(300, 300, 'green'))
    slide3.shapes.add_picture(img_medium, Inches(1), Inches(1), width=Inches(2))

    # Save presentation
    pptx_path = temp_dir / "ranked.pptx"
    prs.save(str(pptx_path))

    # Analyze
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))

    # Should be sorted by size descending
    assert len(results) == 3
    assert results[0]['slide_title'] == "Large Image"
    assert results[1]['slide_title'] == "Medium Image"
    assert results[2]['slide_title'] == "Small Image"

    # Sizes should be descending
    assert results[0]['total_media_bytes'] > results[1]['total_media_bytes']
    assert results[1]['total_media_bytes'] > results[2]['total_media_bytes']


def test_file_not_found():
    """Test that FileNotFoundError is raised for non-existent file."""
    with pytest.raises(FileNotFoundError, match="file not found"):
        pptx_heavy_slides.analyze_pptx_media("/nonexistent/path.pptx")


def test_wrong_file_type(temp_dir):
    """Test that ValueError is raised for non-pptx file."""
    # Create a text file with .txt extension
    txt_path = temp_dir / "test.txt"
    txt_path.write_text("not a presentation")

    with pytest.raises(ValueError, match="unsupported file type"):
        pptx_heavy_slides.analyze_pptx_media(str(txt_path))


def test_corrupt_pptx(temp_dir):
    """Test that ValueError is raised for corrupt .pptx file."""
    # Create a file with .pptx extension but invalid content
    corrupt_path = temp_dir / "corrupt.pptx"
    corrupt_path.write_bytes(b"This is not a valid PPTX file")

    with pytest.raises(ValueError, match="failed to open .pptx"):
        pptx_heavy_slides.analyze_pptx_media(str(corrupt_path))


def test_format_bytes():
    """Test the byte formatting function."""
    assert pptx_heavy_slides.format_bytes(0) == "0.0 MB"
    assert pptx_heavy_slides.format_bytes(500) == "500 B"
    assert pptx_heavy_slides.format_bytes(1024) == "1.0 KB"
    assert pptx_heavy_slides.format_bytes(1024 * 1024) == "1.0 MB"
    assert pptx_heavy_slides.format_bytes(1536 * 1024) == "1.5 MB"
    assert pptx_heavy_slides.format_bytes(1024 * 1024 * 1024) == "1.0 GB"


def test_json_output(temp_dir):
    """Test JSON output generation."""
    # Create a simple presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = "Test"
    img = io.BytesIO(create_test_image(100, 100, 'red'))
    slide.shapes.add_picture(img, Inches(1), Inches(1), width=Inches(2))

    pptx_path = temp_dir / "test.pptx"
    prs.save(str(pptx_path))

    # Analyze and write JSON
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))
    json_path = temp_dir / "output.json"
    pptx_heavy_slides.write_json_output(results, str(json_path))

    # Verify JSON file
    assert json_path.exists()
    with open(json_path) as f:
        data = json.load(f)

    assert isinstance(data, list)
    assert len(data) == 1
    assert data[0]['slide_index'] == 1
    assert data[0]['slide_title'] == "Test"


def test_csv_output(temp_dir):
    """Test CSV output generation."""
    # Create a simple presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = "CSV Test"
    img = io.BytesIO(create_test_image(100, 100, 'blue'))
    slide.shapes.add_picture(img, Inches(1), Inches(1), width=Inches(2))

    pptx_path = temp_dir / "test.pptx"
    prs.save(str(pptx_path))

    # Analyze and write CSV
    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))
    csv_path = temp_dir / "output.csv"
    pptx_heavy_slides.write_csv_output(results, str(csv_path))

    # Verify CSV file
    assert csv_path.exists()
    content = csv_path.read_text()
    lines = content.strip().split('\n')

    assert len(lines) == 2  # Header + 1 data row
    assert 'rank' in lines[0]
    assert 'slide_index' in lines[0]
    assert 'slide_title' in lines[0]
    assert 'CSV Test' in lines[1]


def test_slide_without_title(temp_dir):
    """Test that slides without titles are handled correctly."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout, no title

    img = io.BytesIO(create_test_image(100, 100, 'purple'))
    slide.shapes.add_picture(img, Inches(1), Inches(1), width=Inches(2))

    pptx_path = temp_dir / "no_title.pptx"
    prs.save(str(pptx_path))

    results = pptx_heavy_slides.analyze_pptx_media(str(pptx_path))

    assert len(results) == 1
    assert results[0]['slide_title'] is None
