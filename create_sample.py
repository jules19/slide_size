#!/usr/bin/env python3
"""Create a sample PPTX file for testing the analyzer."""

import io
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def create_image(width, height, color):
    """Create a simple colored image."""
    img = Image.new('RGB', (width, height), color=color)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    return img_bytes

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Large image (will be heaviest)
slide1 = prs.slides.add_slide(prs.slide_layouts[5])
if slide1.shapes.title:
    slide1.shapes.title.text = "Marketing Campaign Photos"
large_img = create_image(1920, 1080, 'blue')
slide1.shapes.add_picture(large_img, Inches(1), Inches(1), width=Inches(6))

# Slide 2: Medium image
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
if slide2.shapes.title:
    slide2.shapes.title.text = "Product Screenshots"
medium_img = create_image(800, 600, 'green')
slide2.shapes.add_picture(medium_img, Inches(2), Inches(2), width=Inches(4))

# Slide 3: Small image (company logo - will appear on multiple slides)
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
if slide3.shapes.title:
    slide3.shapes.title.text = "Company Overview"
small_img = create_image(200, 200, 'red')
slide3.shapes.add_picture(small_img, Inches(3), Inches(3), width=Inches(1.5))

# Slide 4: Same small image (shared media test)
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
if slide4.shapes.title:
    slide4.shapes.title.text = "Contact Information"
small_img2 = create_image(200, 200, 'red')  # Same content as slide 3
slide4.shapes.add_picture(small_img2, Inches(8), Inches(0.5), width=Inches(1))

# Slide 5: No media (text only)
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
if slide5.shapes.title:
    slide5.shapes.title.text = "Thank You"

# Save
prs.save('sample_presentation.pptx')
print("Created sample_presentation.pptx with 5 slides")
